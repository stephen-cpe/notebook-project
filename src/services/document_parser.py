"""Document parser — text extraction from PDF/DOCX/PPTX/TXT/MD.

Each parser reads a file path and returns extracted text. The dispatcher
``extract_text(path, content_type)`` routes by type. When text extraction
yields little/no text (e.g. scanned PDFs), the ingestion pipeline falls back
to GLM-OCR (see ``ocr_service.py``).

Supported types and libraries:
- pdf:  ``pypdf``
- docx: ``python-docx``
- pptx: ``python-pptx``
- txt/md: plain read
"""

from __future__ import annotations

import logging
from pathlib import Path
from typing import Any

from src.services.exceptions import IngestionError

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md"}
TYPE_TO_EXTENSIONS = {
    "pdf": ".pdf",
    "docx": ".docx",
    "pptx": ".pptx",
    "txt": ".txt",
    "md": ".md",
}


def detect_content_type(filename: str) -> str:
    """Return the content type string from a filename's extension.

    Raises:
        IngestionError: if the extension is unsupported or missing.
    """
    ext = Path(filename).suffix.lower()
    for ctype, mapped_ext in TYPE_TO_EXTENSIONS.items():
        if ext == mapped_ext:
            return ctype
    if not ext:
        raise IngestionError(f"Cannot detect content type: {filename!r} has no extension")
    raise IngestionError(f"Unsupported file type: {ext!r} (allowed: {SUPPORTED_EXTENSIONS})")


def extract_text(path: str, content_type: str) -> str:
    """Dispatch to the right parser by ``content_type``.

    Raises:
        IngestionError: if the file is missing or the type is unsupported.
    """
    p = Path(path)
    if not p.exists():
        raise IngestionError(f"File not found: {path}")
    if content_type == "pdf":
        return parse_pdf(path)
    if content_type == "docx":
        return parse_docx(path)
    if content_type == "pptx":
        return parse_pptx(path)
    if content_type in ("txt", "md"):
        return parse_text_file(path)
    raise IngestionError(f"Unsupported content type: {content_type!r}")


def parse_pdf(path: str) -> str:
    """Extract text from a PDF using pypdf (multi-page)."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)
    return "\n\n".join(parts)


def parse_pdf_with_pages(path: str) -> tuple[str, int]:
    """Like ``parse_pdf`` but also returns the page count."""
    from pypdf import PdfReader

    reader = PdfReader(path)
    parts: list[str] = []
    for page in reader.pages:
        text = page.extract_text() or ""
        parts.append(text)
    return "\n\n".join(parts), len(reader.pages)


def parse_docx(path: str) -> str:
    """Extract text from a DOCX using python-docx."""
    from docx import Document

    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text and para.text.strip():
            parts.append(para.text)
    return "\n\n".join(parts)


def parse_pptx(path: str) -> str:
    """Extract text from a PPTX using python-pptx."""
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        slide_texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    txt = para.text
                    if txt and txt.strip():
                        slide_texts.append(txt)
        if slide_texts:
            parts.append("\n".join(slide_texts))
    return "\n\n".join(parts)


def parse_text_file(path: str) -> str:
    """Read a plain text or markdown file as UTF-8."""
    return Path(path).read_text(encoding="utf-8", errors="replace")


def extract_docx_images(path: str) -> list[Any]:
    """Extract all embedded images from a DOCX as PIL Images.

    DOCX files are ZIP archives; embedded images live under ``word/media/``.
    Returns a list of PIL Images. Returns an empty list if the file contains
    no images or extraction fails for any reason (graceful degradation, FR-24).
    """
    return _extract_zip_media(path, "word/media/")


def extract_pptx_images(path: str) -> list[Any]:
    """Extract all embedded images from a PPTX as PIL Images.

    PPTX files are ZIP archives; embedded images live under ``ppt/media/``.
    Returns a list of PIL Images. Returns an empty list if the file contains
    no images or extraction fails for any reason (graceful degradation, FR-24).
    """
    return _extract_zip_media(path, "ppt/media/")


def _extract_zip_media(path: str, media_prefix: str) -> list[Any]:
    """Extract images from a ZIP-based Office file under ``media_prefix``.

    Used by both DOCX (``word/media/``) and PPTX (``ppt/media/``) extraction.
    """
    import io
    import zipfile

    from PIL import Image

    images: list[Any] = []
    try:
        with zipfile.ZipFile(path) as z:
            for name in z.namelist():
                if not name.startswith(media_prefix):
                    continue
                with z.open(name) as fh:
                    images.append(Image.open(io.BytesIO(fh.read())).convert("RGB"))
    except Exception as exc:  # noqa: BLE001
        logger.warning("Failed to extract images from %s: %s", path, exc)
        return []
    return images


# Magic bytes for supported file types (NFR-23).
_MAGIC_BYTES: dict[str, bytes] = {
    "pdf": b"%PDF",
    "docx": b"PK\x03\x04",
    "pptx": b"PK\x03\x04",
    # txt/md have no magic bytes — validated as plain text (no null bytes).
}


def validate_magic_bytes(file_path: str, content_type: str) -> bool:
    """Check that a file's first bytes match the expected magic bytes for its type.

    For txt/md, verifies the file contains no null bytes (plain text check).
    Returns True if the file passes validation, False otherwise.
    """
    p = Path(file_path)
    if not p.exists():
        return False

    if content_type in ("txt", "md"):
        # Plain text: read first 4 KB and check for null bytes.
        with p.open("rb") as f:
            chunk = f.read(4096)
        return b"\x00" not in chunk

    expected = _MAGIC_BYTES.get(content_type)
    if expected is None:
        return False

    with p.open("rb") as f:
        header = f.read(len(expected))
    return header == expected
