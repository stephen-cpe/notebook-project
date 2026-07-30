"""Generate test fixture files for the document parser tests.

Creates: sample.pdf, sample.docx, sample.pptx, sample.txt, sample.md
under tests/fixtures/. Run with: ``python tests/fixtures/_generate.py``
"""

from __future__ import annotations

from pathlib import Path

FIXTURES = Path(__file__).resolve().parent


def make_pdf() -> None:
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", size=12)
    pdf.cell(200, 10, text="PDF Fixture Title Page", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(10)
    pdf.multi_cell(
        0,
        8,
        text=(
            "This is page one of the PDF fixture. "
            "It contains searchable text about machine learning and data science. "
            "The quick brown fox jumps over the lazy dog."
        ),
    )
    pdf.add_page()
    pdf.multi_cell(
        0,
        8,
        text=(
            "This is page two. It discusses neural networks and deep learning. "
            "Convolutional layers extract spatial features from images."
        ),
    )
    pdf.output(str(FIXTURES / "sample.pdf"))


def make_docx() -> None:
    from docx import Document

    doc = Document()
    doc.add_heading("DOCX Fixture Document", level=1)
    doc.add_paragraph(
        "This is a DOCX fixture document. It contains text about "
        "artificial intelligence and natural language processing."
    )
    doc.add_paragraph(
        "A second paragraph covers transformers, attention mechanisms, "
        "and the history of large language models."
    )
    doc.save(str(FIXTURES / "sample.docx"))


def make_pptx() -> None:
    from pptx import Presentation

    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    shapes = slide.shapes
    shapes.title.text = "PPTX Fixture Slide One"
    shapes.placeholders[1].text = "Slide one text about cloud computing and distributed systems."
    slide2 = prs.slides.add_slide(slide_layout)
    shapes2 = slide2.shapes
    shapes2.title.text = "PPTX Fixture Slide Two"
    shapes2.placeholders[1].text = "Slide two text about Kubernetes and containers."
    prs.save(str(FIXTURES / "sample.pptx"))


def make_txt() -> None:
    (FIXTURES / "sample.txt").write_text(
        "TXT fixture document.\n"
        "Plain text about databases, SQL, and NoSQL stores.\n"
        "PostgreSQL, MongoDB, and Redis are common choices.\n",
        encoding="utf-8",
    )


def make_md() -> None:
    (FIXTURES / "sample.md").write_text(
        "# Markdown Fixture\n\n"
        "A markdown document about web development.\n\n"
        "## Frontend\n\n"
        "React, Vue, and Svelte are popular frameworks.\n\n"
        "## Backend\n\n"
        "Flask and FastAPI are Python web frameworks.\n",
        encoding="utf-8",
    )


def make_empty_pdf() -> None:
    """A PDF with no selectable text (blank page) for OCR fallback tests."""
    from fpdf import FPDF

    pdf = FPDF()
    pdf.add_page()
    pdf.output(str(FIXTURES / "empty.pdf"))


if __name__ == "__main__":
    make_pdf()
    make_docx()
    make_pptx()
    make_txt()
    make_md()
    make_empty_pdf()
    print("Fixtures generated in", FIXTURES)
